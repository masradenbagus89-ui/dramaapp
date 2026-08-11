# -*- coding: utf-8 -*-
"""Generator PowerPoint — Gambaran Teknis Project DramaApp / DramaKu."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------
# Palet warna (mengikuti tema gelap-sinematik aplikasi)
# ----------------------------------------------------------------------
BG      = RGBColor(0x0B, 0x0B, 0x12)
BG_ALT  = RGBColor(0x0E, 0x0A, 0x16)
CARD    = RGBColor(0x18, 0x18, 0x26)
CARD2   = RGBColor(0x21, 0x16, 0x2C)
LINEC   = RGBColor(0x2C, 0x2C, 0x3D)
ROSE    = RGBColor(0xF4, 0x3F, 0x5E)
PINK    = RGBColor(0xEC, 0x48, 0x99)
PURPLE  = RGBColor(0xA8, 0x55, 0xF7)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xDD, 0xDF, 0xE6)
GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DIM     = RGBColor(0x60, 0x64, 0x74)
GLOW_R  = RGBColor(0x33, 0x12, 0x20)
GLOW_P  = RGBColor(0x23, 0x16, 0x38)

FONT  = "Segoe UI"
MONO  = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# Total slide & penomoran otomatis — nomor mengikuti urutan new_slide()
TOTAL = 18
CUR = 0


# ----------------------------------------------------------------------
# Helper
# ----------------------------------------------------------------------
def new_slide(bg=BG):
    global CUR
    CUR += 1
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def notes(slide, text):
    """Naskah pembicara (speaker notes) — muncul di panel Notes PowerPoint."""
    slide.notes_slide.notes_text_frame.text = text.strip()


def rect(slide, kind, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def grad(slide, kind, x, y, w, h, c1, c2, angle=0, radius=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.line.fill.background()
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def txt(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """paras: list of paragraphs; tiap paragraph list of run-tuple
       (text, size, color, bold[, font, italic])."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for run in para:
            t, size, color, bold = run[0], run[1], run[2], run[3]
            fn = run[4] if len(run) > 4 else FONT
            it = run[5] if len(run) > 5 else False
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = fn
            r.font.italic = it
    return tb


def line(slide, x, y, w, h=0.022, color=LINEC):
    rect(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=color)


def header(slide, kicker, title, num=None):
    num = CUR  # penomoran otomatis mengikuti urutan slide
    bar = rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 0.66, 0.13, 0.62,
               fill=ROSE, radius=0.5)
    txt(slide, 0.95, 0.6, 9.5, 0.32, [[(kicker.upper(), 12, PURPLE, True)]])
    txt(slide, 0.93, 0.86, 11.0, 0.7, [[(title, 29, WHITE, True)]])
    txt(slide, 11.4, 0.66, 1.3, 0.5,
        [[(f"{num:02d}", 14, DIM, True)]], align=PP_ALIGN.RIGHT)
    txt(slide, 11.4, 0.92, 1.3, 0.4,
        [[(f"/ {TOTAL}", 11, DIM, False)]], align=PP_ALIGN.RIGHT)
    line(slide, 0.95, 1.62, 11.45)


def footer(slide, n=None):
    n = CUR  # penomoran otomatis
    txt(slide, 0.95, 7.02, 9.0, 0.32,
        [[("DramaKu", 9, ROSE, True), ("  ·  Gambaran Teknis Project", 9, DIM, False)]])
    txt(slide, 11.4, 7.02, 1.3, 0.32,
        [[(f"{n:02d}", 9, DIM, True)]], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=CARD, radius=0.055):
    return rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, radius=radius)


def chip(slide, x, y, w, label, color=LIGHT, fill=CARD, size=11, h=0.44):
    c = rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, radius=0.5)
    c.line.color.rgb = LINEC
    c.line.width = Pt(0.75)
    txt(slide, x, y, w, h, [[(label, size, color, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 1 — Judul
# ======================================================================
s = new_slide(BG)
# glow dekoratif
g = rect(s, MSO_SHAPE.OVAL, 8.6, -2.4, 6.6, 6.6, fill=GLOW_P)
g = rect(s, MSO_SHAPE.OVAL, 10.4, 2.6, 5.4, 5.4, fill=GLOW_R)
rect(s, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, fill=ROSE)

grad(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.05, 1.18, 1.7, 0.5, ROSE, PURPLE,
     angle=0, radius=0.5)
txt(s, 1.05, 1.18, 1.7, 0.5, [[("PROJECT", 11, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

txt(s, 1.0, 2.0, 11.0, 1.7,
    [[("Drama", 88, WHITE, True), ("Ku", 88, ROSE, True)]])
txt(s, 1.07, 3.5, 11.0, 0.6,
    [[("Platform Streaming Drama China Pendek", 23, LIGHT, False)]])
txt(s, 1.07, 4.06, 11.2, 0.5,
    [[("Aplikasi web modern — mobile-first, gratis, ringan", 14, GRAY, False)]])

grad(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.07, 4.78, 3.0, 0.07, ROSE, PURPLE, angle=0)

chip(s, 1.07, 5.1, 1.95, "Next.js 16", PURPLE)
chip(s, 3.17, 5.1, 1.7, "React 19", PURPLE)
chip(s, 5.02, 5.1, 2.05, "TypeScript", PURPLE)
chip(s, 7.22, 5.1, 1.95, "Tailwind 4", PURPLE)

txt(s, 1.07, 6.35, 11.0, 0.4,
    [[("GAMBARAN TEKNIS PROJECT", 12, ROSE, True)]])
txt(s, 1.07, 6.68, 11.0, 0.4,
    [[("Presentasi  ·  2 Juni 2026  ·  github.com/masradenbagus89-ui/dramaapp",
       11, DIM, False)]])
notes(s, """
PEMBUKA. Perkenalkan diri lalu sampaikan kalimat ini:
"DramaApp (brand: DramaKu) adalah platform streaming drama China pendek yang
saya bangun dari nol — sudah ONLINE dan bisa diakses siapa saja di
dramaapp.vercel.app. Saya akan jelaskan dari dua sisi: sisi teknis aplikasinya,
dan sistem yang saya pakai selama membangunnya."
TIPS: kalau memungkinkan, buka dramaapp.vercel.app langsung sebagai demo singkat
sebelum lanjut — bukti app benar-benar jalan paling meyakinkan.
Teknologi inti: Next.js 16, React 19, TypeScript, Tailwind 4. Biaya hosting Rp 0.
""")

# ======================================================================
# SLIDE 2 — Agenda
# ======================================================================
s = new_slide(BG)
header(s, "Daftar Isi", "Apa yang akan dibahas", 2)
agenda = [
    ("01", "Gambaran Produk", "Masalah yang dipecahkan, fitur utama bagi pengguna", ROSE),
    ("02", "Tech Stack & Arsitektur", "Teknologi inti dan cara komponen disusun", PURPLE),
    ("03", "Routing, API & Data", "Struktur halaman, backend, dan model data", AMBER),
    ("04", "Admin, Video & Deployment", "Panel admin, streaming video, rilis ke cloud", TEAL),
]
cw, gap = 2.69, 0.27
x0 = 0.95
for i, (no, t, d, col) in enumerate(agenda):
    x = x0 + i * (cw + gap)
    card(s, x, 2.1, cw, 4.4)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, cw, 0.12, fill=col, radius=0.5)
    txt(s, x + 0.28, 2.5, cw - 0.5, 0.9, [[(no, 40, col, True)]])
    txt(s, x + 0.28, 3.45, cw - 0.56, 1.0,
        [[(t, 16, WHITE, True)]], line_spacing=1.05)
    line(s, x + 0.28, 4.35, cw - 0.56)
    txt(s, x + 0.28, 4.55, cw - 0.56, 1.7,
        [[(d, 12, GRAY, False)]], line_spacing=1.3)
footer(s, 2)
notes(s, """
AGENDA. "Presentasi dibagi 4 bagian: (1) gambaran produk & fitur, (2) tech stack
dan arsitektur, (3) routing, API, dan data, (4) admin, video, deployment — plus
arsitektur hybrid, perjalanan pembuatan, dan tanya-jawab di akhir."
Tidak perlu dibaca satu-satu; cukup sebut empat bagian besar lalu lanjut.
""")

# ======================================================================
# SLIDE 3 — Apa itu DramaApp
# ======================================================================
s = new_slide(BG)
header(s, "Gambaran Produk", "Apa itu DramaApp?", 3)
txt(s, 0.95, 1.95, 7.1, 1.3,
    [[("Aplikasi web untuk ", 16, LIGHT, False),
      ("streaming drama China pendek", 16, ROSE, True),
      (" secara gratis — dirancang ", 16, LIGHT, False),
      ("mobile-first", 16, AMBER, True),
      (" agar nyaman ditonton langsung dari HP maupun browser desktop.",
       16, LIGHT, False)]],
    line_spacing=1.35)

items = [
    ("Tren short drama", "Drama China format pendek sedang populer — butuh platform yang ringan & cepat."),
    ("Akses gratis & mudah", "Tidak perlu instal aplikasi — cukup buka di browser HP atau PC."),
    ("Pengalaman terkurasi", "Katalog rapi per kategori, halaman detail, dan feed Shorts vertikal."),
    ("Bisa dikelola sendiri", "Panel admin internal untuk menambah judul, episode, dan video."),
]
y = 3.3
for head, body in items:
    rect(s, MSO_SHAPE.OVAL, 0.97, y + 0.07, 0.16, 0.16, fill=ROSE)
    txt(s, 1.32, y, 6.8, 0.9,
        [[(head + "  —  ", 13.5, WHITE, True), (body, 13.5, GRAY, False)]],
        line_spacing=1.25)
    y += 0.84

# panel kanan — statistik
px = 8.45
card(s, px, 1.95, 3.95, 4.55, fill=CARD2)
txt(s, px, 2.2, 3.95, 0.4,
    [[("DRAMAKU — DRAMA CHINA GRATIS", 10.5, PURPLE, True)]], align=PP_ALIGN.CENTER)
line(s, px + 0.4, 2.62, 3.15)
stats = [("33", "judul drama dalam katalog", ROSE),
         ("7", "kategori genre", AMBER),
         ("25", "poster siap tampil", PURPLE),
         ("11", "halaman / rute aplikasi", TEAL)]
sy = 2.8
for big, lbl, col in stats:
    txt(s, px + 0.4, sy, 1.2, 0.7, [[(big, 30, col, True)]])
    txt(s, px + 1.65, sy + 0.12, 2.1, 0.7,
        [[(lbl, 11.5, LIGHT, False)]], line_spacing=1.1,
        anchor=MSO_ANCHOR.MIDDLE)
    sy += 0.92
footer(s, 3)
notes(s, """
APA ITU DRAMAAPP. Inti: aplikasi web streaming drama China pendek, GRATIS,
dirancang mobile-first (utamanya untuk HP). Konsepnya gabungan Netflix (katalog
rapi per kategori + halaman detail) dan TikTok (feed Shorts vertikal).
Kenapa dibuat: tren short drama sedang naik, butuh platform ringan & cepat,
tanpa perlu instal aplikasi — cukup buka browser.
Angka untuk diingat: 33 judul drama, 7 kategori genre, ~25 poster, 11 halaman.
""")

# ======================================================================
# SLIDE 4 — Fitur Utama
# ======================================================================
s = new_slide(BG)
header(s, "Gambaran Produk", "Fitur Utama Aplikasi", 4)
feats = [
    ("Katalog & Discover", "Jelajahi semua drama, filter per kategori, cari judul.", ROSE),
    ("Streaming per Episode", "Pemutar video khusus tiap episode dengan navigasi prev/next.", PURPLE),
    ("Feed Shorts", "Tampilan potongan drama bergaya feed vertikal.", AMBER),
    ("Akun & Profil", "Login / daftar, peran admin & viewer, avatar berwarna.", TEAL),
    ("Interaksi Sosial", "Suka (like), komentar, dan simpan ke My List.", PINK),
    ("Panel Admin", "Kelola judul, episode, poster, dan berkas video.", PURPLE),
]
cw, ch, gx, gy = 3.69, 1.95, 0.19, 0.22
x0, y0 = 0.95, 2.0
for i, (t, d, col) in enumerate(feats):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    card(s, x, y, cw, ch)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.32, y + 0.34, 0.5, 0.11,
         fill=col, radius=0.5)
    txt(s, x + 0.32, y + 0.58, cw - 0.6, 0.5, [[(t, 15, WHITE, True)]])
    txt(s, x + 0.32, y + 1.0, cw - 0.62, 0.85,
        [[(d, 11.5, GRAY, False)]], line_spacing=1.25)
footer(s, 4)
notes(s, """
FITUR UTAMA. Enam fitur inti dari sisi pengguna:
1. Katalog & Discover — jelajah, filter per kategori.
2. Streaming per episode — player khusus tiap episode, navigasi prev/next.
3. Feed Shorts — scroll vertikal ala TikTok.
4. Akun & profil — login/daftar, role admin & viewer, avatar warna.
5. Interaksi sosial — like, komentar, simpan ke My List.
6. Panel admin — kelola judul, episode, poster, video.
Detail menarik: VideoPlayer punya fallback otomatis — kalau video asli belum
ada, tampil video sample + label, jadi app tetap jalan.
""")

# ======================================================================
# SLIDE 5 — Tech Stack
# ======================================================================
s = new_slide(BG)
header(s, "Fondasi Teknis", "Tech Stack", 5)
groups = [
    ("FRONTEND", ROSE, [
        ("Next.js 16.2.4", "Framework React, App Router"),
        ("React 19.2.4", "Library UI komponen"),
        ("TypeScript 5", "Bahasa dengan tipe statis"),
    ]),
    ("STYLING & UI", PURPLE, [
        ("Tailwind CSS 4", "Utility-first styling"),
        ("next/font", "Geist & Playfair Display"),
        ("Desain mobile-first", "Top & Bottom navigation"),
    ]),
    ("BUILD & RUNTIME", AMBER, [
        ("Turbopack", "Bundler dev super cepat"),
        ("Node.js 24", "Runtime server-side"),
        ("Route Handlers", "API bawaan Next.js"),
    ]),
    ("DELIVERY", TEAL, [
        ("Vercel", "Hosting & deploy otomatis"),
        ("GitHub", "Repositori publik"),
        ("Cloudflare R2", "Penyimpanan video (rencana)"),
    ]),
]
cw, gap = 2.69, 0.27
x0 = 0.95
for i, (gt, col, rows) in enumerate(groups):
    x = x0 + i * (cw + gap)
    card(s, x, 2.0, cw, 4.55)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.0, cw, 0.12, fill=col, radius=0.5)
    txt(s, x + 0.26, 2.32, cw - 0.5, 0.4, [[(gt, 12, col, True)]])
    line(s, x + 0.26, 2.74, cw - 0.52)
    ry = 2.95
    for name, desc in rows:
        txt(s, x + 0.26, ry, cw - 0.5, 0.4, [[(name, 13, WHITE, True)]])
        txt(s, x + 0.26, ry + 0.3, cw - 0.5, 0.6,
            [[(desc, 10.5, GRAY, False)]], line_spacing=1.15)
        ry += 1.13
footer(s, 5)
notes(s, """
TECH STACK. Jelaskan per kelompok:
- FRONTEND: Next.js 16 (framework React, App Router), React 19, TypeScript 5
  (tipe statis = lebih aman dari error).
- STYLING: Tailwind CSS 4 (desain cepat & responsive), next/font, mobile-first.
- BUILD/RUNTIME: Turbopack (bundler super cepat), Node.js 24, Route Handlers
  (API bawaan Next.js — backend & frontend satu kerangka).
- DELIVERY: Vercel (hosting + deploy otomatis dari GitHub), GitHub (repo publik),
  Cloudflare R2 untuk penyimpanan video (rencana).
Pesan kunci: ini FULL-STACK dalam satu codebase, teknologinya modern & standar
industri.
""")

# ======================================================================
# SLIDE 6 — Arsitektur Aplikasi
# ======================================================================
s = new_slide(BG)
header(s, "Fondasi Teknis", "Arsitektur Aplikasi", 6)
txt(s, 0.95, 1.9, 6.7, 0.9,
    [[("Aplikasi penuh (full-stack) dalam satu codebase Next.js — "
       "lapisan tampilan dan backend hidup berdampingan.", 13.5, LIGHT, False)]],
    line_spacing=1.3)

arch = [
    ("Server Components", "Halaman dirender di server secara default — cepat & ringan di sisi klien.", ROSE),
    ("Client Components", "Bagian interaktif (\"use client\") seperti pemutar video, tombol like, komentar.", PURPLE),
    ("Route Handlers (API)", "Endpoint backend di app/api/* — menyajikan & menyimpan data.", AMBER),
    ("Penyimpanan berbasis file", "Data disimpan sebagai berkas JSON — tanpa database, dibaca/tulis lewat Node fs.", TEAL),
    ("State di klien", "Status login, like, dan My List disimpan di localStorage browser.", PINK),
]
y = 2.95
for head, body, col in arch:
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.97, y + 0.03, 0.13, 0.5,
         fill=col, radius=0.5)
    txt(s, 1.28, y, 6.3, 0.4, [[(head, 13.5, WHITE, True)]])
    txt(s, 1.28, y + 0.3, 6.35, 0.55,
        [[(body, 11, GRAY, False)]], line_spacing=1.18)
    y += 0.78

# diagram stack di kanan
px = 8.2
txt(s, px, 1.9, 4.2, 0.35, [[("ALUR LAPISAN", 11, PURPLE, True)]])
layers = [
    ("KLIEN", "Browser · localStorage · Client Components", ROSE),
    ("RENDERING", "React Server Components (App Router)", PURPLE),
    ("API", "Route Handlers  —  /api/*", AMBER),
    ("DATA", "Berkas JSON  +  Video CDN (R2)", TEAL),
]
ly = 2.35
for i, (t, d, col) in enumerate(layers):
    card(s, px, ly, 4.2, 0.92, fill=CARD2)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, px, ly, 0.13, 0.92, fill=col, radius=0.5)
    txt(s, px + 0.32, ly + 0.13, 3.8, 0.35, [[(t, 12.5, col, True)]])
    txt(s, px + 0.32, ly + 0.45, 3.8, 0.4, [[(d, 10, LIGHT, False)]])
    if i < len(layers) - 1:
        txt(s, px + 1.9, ly + 0.93, 0.5, 0.3,
            [[("▼", 10, DIM, False)]], align=PP_ALIGN.CENTER)
    ly += 1.12
footer(s, 6)
notes(s, """
ARSITEKTUR APLIKASI. Satu codebase Next.js berisi frontend + backend.
- Server Components: halaman dirender di server (cepat, ringan di HP).
- Client Components ("use client"): bagian interaktif — player video, tombol like,
  komentar.
- Route Handlers (app/api/*): endpoint backend untuk baca & simpan data.
- Penyimpanan berbasis file: data sebagai JSON, dibaca/tulis via Node fs — TANPA
  database.
- State di klien: status login, like, My List disimpan di localStorage browser.
Diagram kanan menunjukkan alur lapisan: Klien → Rendering → API → Data.
""")

# ======================================================================
# SLIDE 7 — Routing / Struktur Halaman
# ======================================================================
s = new_slide(BG)
header(s, "Routing, API & Data", "Struktur Halaman & Routing", 7)
txt(s, 0.95, 1.85, 11.4, 0.4,
    [[("11 rute di App Router — penamaan folder = penamaan URL.",
       12.5, GRAY, False)]])

routes_l = [
    ("/", "Landing page — pintu masuk aplikasi"),
    ("/login", "Masuk ke akun"),
    ("/daftar", "Registrasi akun baru"),
    ("/beranda", "Beranda — daftar drama utama"),
    ("/discover", "Jelajah & filter per kategori"),
    ("/shorts", "Feed potongan drama vertikal"),
]
routes_r = [
    ("/drama/[id]", "Halaman detail tiap drama"),
    ("/watch/[id]/[ep]", "Pemutar video per episode"),
    ("/my-list", "Daftar drama yang disimpan"),
    ("/profile", "Profil & pengaturan pengguna"),
    ("/admin", "Panel pengelolaan konten"),
]

def route_col(slide, x, y, items):
    cy = y
    for route, desc in items:
        card(slide, x, cy, 5.55, 0.72)
        rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, 0.1, 0.72,
             fill=AMBER, radius=0.5)
        txt(slide, x + 0.3, cy + 0.1, 5.1, 0.35,
            [[(route, 13, AMBER, True, MONO)]])
        txt(slide, x + 0.3, cy + 0.4, 5.1, 0.3,
            [[(desc, 10.5, GRAY, False)]])
        cy += 0.83

route_col(s, 0.95, 2.35, routes_l)
route_col(s, 6.83, 2.35, routes_r)
txt(s, 6.83, 6.5, 5.55, 0.4,
    [[("[id] & [ep] = segmen dinamis — satu file melayani banyak URL.",
       10.5, DIM, False, FONT, True)]])
footer(s, 7)
notes(s, """
STRUKTUR HALAMAN & ROUTING. Next.js App Router: nama folder = nama URL.
Ada 11 rute. Untuk viewer: /, /login, /daftar, /beranda, /discover, /shorts,
/drama/[id], /watch/[id]/[ep], /my-list, /profile. Untuk admin: /admin.
Poin penting: [id] dan [ep] = segmen DINAMIS — satu file melayani banyak URL
(mis. semua drama pakai satu file /drama/[id]/page.tsx).
""")

# ======================================================================
# SLIDE 8 — API & Backend
# ======================================================================
s = new_slide(BG)
header(s, "Routing, API & Data", "API & Lapisan Backend", 8)
txt(s, 0.95, 1.85, 11.4, 0.4,
    [[("Route Handlers di app/api/* — sebagian publik, sebagian khusus admin.",
       12.5, GRAY, False)]])

# Publik
card(s, 0.95, 2.35, 5.55, 4.2)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, 2.35, 5.55, 0.12, fill=TEAL, radius=0.5)
txt(s, 1.25, 2.6, 5.0, 0.4, [[("API PUBLIK", 13, TEAL, True)]])
line(s, 1.25, 3.0, 4.95)
pub = [
    ("GET /api/dramas", "Ambil seluruh daftar drama"),
    ("/api/comments", "Baca & kirim komentar"),
    ("/api/likes", "Catat suka pada drama"),
    ("/api/admins", "Daftar email admin terdaftar"),
]
py = 3.15
for ep, d in pub:
    txt(s, 1.25, py, 5.0, 0.35, [[(ep, 12.5, AMBER, True, MONO)]])
    txt(s, 1.25, py + 0.28, 5.0, 0.3, [[(d, 10.5, GRAY, False)]])
    py += 0.8

# Admin
card(s, 6.83, 2.35, 5.55, 4.2, fill=CARD2)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.83, 2.35, 5.55, 0.12, fill=ROSE, radius=0.5)
txt(s, 7.13, 2.6, 5.0, 0.4, [[("API ADMIN  (terproteksi)", 13, ROSE, True)]])
line(s, 7.13, 3.0, 4.95)
adm = [
    ("/api/admin/drama", "Tambah, ubah, hapus drama (CRUD)"),
    ("/api/admin/scan", "Pindai berkas video tiap drama"),
    ("/api/admin/upload", "Unggah poster & aset"),
    ("/api/admin/hardlink", "Tautkan berkas video ke aplikasi"),
]
py = 3.15
for ep, d in adm:
    txt(s, 7.13, py, 5.0, 0.35, [[(ep, 12.5, AMBER, True, MONO)]])
    txt(s, 7.13, py + 0.28, 5.0, 0.3, [[(d, 10.5, GRAY, False)]])
    py += 0.8

txt(s, 7.13, 6.05, 5.0, 0.45,
    [[("Proteksi: header ", 10, GRAY, False),
      ("x-admin-email", 10, AMBER, True, MONO),
      (" divalidasi vs admins.json", 10, GRAY, False)]])
footer(s, 8)
notes(s, """
API & BACKEND. Endpoint ada di app/api/* — sebagian publik, sebagian admin.
PUBLIK: GET /api/dramas (daftar drama), /api/comments (komentar),
/api/likes (like), /api/admins (daftar email admin).
ADMIN (terproteksi): /api/admin/drama (CRUD drama), /api/admin/scan (pindai
video), /api/admin/upload (poster/aset), /api/admin/hardlink (tautkan video).
KEAMANAN: tiap request admin membawa header x-admin-email yang dicocokkan ke
admins.json; kalau tak cocok → ditolak 401. Validasi di SERVER, bukan tampilan.
""")

# ======================================================================
# SLIDE 9 — Model Data
# ======================================================================
s = new_slide(BG)
header(s, "Routing, API & Data", "Model Data", 9)
txt(s, 0.95, 1.9, 6.6, 0.85,
    [[("Tanpa database — seluruh data tersimpan sebagai berkas ", 13, LIGHT, False),
      ("JSON", 13, AMBER, True),
      (" di folder ", 13, LIGHT, False),
      ("data/", 13, AMBER, True, MONO),
      (".", 13, LIGHT, False)]], line_spacing=1.3)

files = [
    ("dramas.json", "33 drama — judul, kategori, episode, sinopsis"),
    ("comments.json", "Komentar pengguna per drama"),
    ("interactions.json", "Catatan suka / interaksi"),
    ("admins.json", "Daftar email yang berperan admin"),
]
y = 2.85
for f, d in files:
    rect(s, MSO_SHAPE.OVAL, 0.99, y + 0.07, 0.15, 0.15, fill=ROSE)
    txt(s, 1.3, y, 6.3, 0.6,
        [[(f, 12.5, AMBER, True, MONO), ("   " + d, 11, GRAY, False)]],
        line_spacing=1.2)
    y += 0.62

txt(s, 0.95, 5.5, 6.6, 0.35, [[("7 KATEGORI GENRE", 11, PURPLE, True)]])
cats = ["Romance", "Tycoon", "Harem", "Time Travel", "Action", "Comedy", "Fantasy"]
cx = 0.95
for c in cats:
    w = 0.28 + len(c) * 0.085
    chip(s, cx, 5.85, w, c, LIGHT, CARD, 9.5, h=0.36)
    cx += w + 0.13

# kartu kode tipe Drama
px = 8.0
card(s, px, 1.95, 4.4, 4.7, fill=RGBColor(0x10, 0x12, 0x1C))
txt(s, px + 0.3, 2.15, 3.9, 0.35, [[("lib/types.ts", 10.5, DIM, True, MONO)]])
line(s, px + 0.3, 2.5, 3.8, color=RGBColor(0x26,0x28,0x36))
code = [
    [("type ", 11.5, PURPLE, False, MONO), ("Drama", 11.5, TEAL, True, MONO),
     (" = {", 11.5, LIGHT, False, MONO)],
    [("  id, title, category,", 11.5, LIGHT, False, MONO)],
    [("  episodes, views,", 11.5, LIGHT, False, MONO)],
    [("  synopsis, gradient,", 11.5, LIGHT, False, MONO)],
    [("  posterImage?,", 11.5, AMBER, False, MONO)],
    [("  heroImage?,", 11.5, AMBER, False, MONO)],
    [("  exclusive?", 11.5, AMBER, False, MONO)],
    [("}", 11.5, LIGHT, False, MONO)],
]
txt(s, px + 0.3, 2.7, 3.9, 2.6, code, line_spacing=1.5)
line(s, px + 0.3, 5.15, 3.8, color=RGBColor(0x26,0x28,0x36))
txt(s, px + 0.3, 5.32, 3.9, 1.2,
    [[("? = field opsional", 10, DIM, False, MONO)],
     [("ID dibuat otomatis dari judul", 10, GRAY, False)],
     [("(slugify) saat drama ditambah.", 10, GRAY, False)]],
    line_spacing=1.3, space_after=3)
footer(s, 9)
notes(s, """
MODEL DATA. Tanpa database — semua data sebagai berkas JSON di folder data/:
dramas.json (33 drama: judul, kategori, episode, sinopsis), comments.json
(komentar per drama), interactions.json (like/interaksi), admins.json (email
admin). Tipe Drama didefinisikan di lib/types.ts; field ber-"?" = opsional.
ID drama dibuat OTOMATIS dari judul (slugify) saat ditambahkan.
7 kategori: Romance, Tycoon, Harem, Time Travel, Action, Comedy, Fantasy.
Kalau ditanya kenapa tanpa DB → lihat slide Tantangan: cukup untuk skala awal,
mudah di-upgrade ke database nanti.
""")

# ======================================================================
# SLIDE 10 — Panel Admin
# ======================================================================
s = new_slide(BG)
header(s, "Admin, Video & Deployment", "Panel Admin", 10)
txt(s, 0.95, 1.9, 11.4, 0.5,
    [[("Halaman ", 13, LIGHT, False), ("/admin", 13, AMBER, True, MONO),
      (" — pengelolaan konten tanpa perlu menyentuh kode.", 13, LIGHT, False)]])

adm_feats = [
    ("Kelola Drama", "Tambah, ubah, dan hapus judul drama (operasi CRUD penuh).", ROSE),
    ("Pindai Video", "Deteksi berkas video yang tersedia untuk tiap drama.", PURPLE),
    ("Unggah Aset", "Upload poster dan gambar pendukung lewat antarmuka.", AMBER),
    ("Tautkan Video", "Hardlink berkas video besar agar siap diputar aplikasi.", TEAL),
]
cw, gx = 2.69, 0.27
x0 = 0.95
for i, (t, d, col) in enumerate(adm_feats):
    x = x0 + i * (cw + gx)
    card(s, x, 2.55, cw, 2.45)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.28, 2.85, 0.55, 0.12,
         fill=col, radius=0.5)
    txt(s, x + 0.28, 3.12, cw - 0.5, 0.45, [[(t, 14.5, WHITE, True)]])
    txt(s, x + 0.28, 3.55, cw - 0.54, 1.3,
        [[(d, 11, GRAY, False)]], line_spacing=1.3)

card(s, 0.95, 5.25, 11.45, 1.25, fill=CARD2)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, 5.25, 0.13, 1.25, fill=ROSE, radius=0.5)
txt(s, 1.3, 5.45, 1.9, 0.4, [[("KEAMANAN", 11, ROSE, True)]])
txt(s, 1.3, 5.78, 10.9, 0.65,
    [[("Setiap permintaan ke API admin memuat header ", 12, LIGHT, False),
      ("x-admin-email", 12, AMBER, True, MONO),
      (" yang dicocokkan dengan daftar di ", 12, LIGHT, False),
      ("admins.json", 12, AMBER, True, MONO),
      (" — bila tidak cocok, ditolak (401 Unauthorized).", 12, LIGHT, False)]],
    line_spacing=1.3)
footer(s, 10)
notes(s, """
PANEL ADMIN. Halaman /admin = kelola konten TANPA menyentuh kode.
Empat fungsi: Kelola drama (CRUD), Pindai video (deteksi file tersedia tiap
drama), Unggah aset (poster/gambar), Tautkan video (hardlink file besar agar
siap diputar). Bawah: pengingat keamanan — header x-admin-email dicocokkan ke
admins.json, kalau gagal 401.
Cerita evolusi (bisa diceritakan di sini): dulu tambah drama = edit file JSON +
copy video manual; sekarang lewat form jadi praktis. Detail "3 klik" ada di
slide Perjalanan Pembuatan.
""")

# ======================================================================
# SLIDE 11 — Video Streaming
# ======================================================================
s = new_slide(BG)
header(s, "Admin, Video & Deployment", "Video Streaming", 11)

txt(s, 0.95, 1.95, 6.7, 0.4, [[("CARA KERJA PEMUTAR", 11, PURPLE, True)]])
vid = [
    ("Komponen VideoPlayer", "Memakai elemen HTML5 <video> bawaan browser."),
    ("Sumber video dinamis", "URL dirakit: {BASE}/{id-drama}/{episode}.mp4."),
    ("Fallback otomatis", "Bila video gagal dimuat, beralih ke /sample.mp4."),
    ("Navigasi episode", "Tombol episode sebelumnya & berikutnya di halaman watch."),
]
y = 2.4
for head, body in vid:
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.97, y + 0.03, 0.13, 0.52,
         fill=ROSE, radius=0.5)
    txt(s, 1.28, y, 6.4, 0.4, [[(head, 13.5, WHITE, True)]])
    txt(s, 1.28, y + 0.3, 6.45, 0.5,
        [[(body, 11, GRAY, False)]], line_spacing=1.2)
    y += 0.92

# kartu env
card(s, 8.1, 2.4, 4.3, 1.5, fill=RGBColor(0x10, 0x12, 0x1C))
txt(s, 8.35, 2.58, 3.9, 0.3, [[("ENV VARIABLE", 9.5, DIM, True, MONO)]])
txt(s, 8.35, 2.86, 3.9, 0.4,
    [[("NEXT_PUBLIC_VIDEO_", 11, TEAL, True, MONO)]])
txt(s, 8.35, 3.08, 3.9, 0.4,
    [[("BASE_URL", 11, TEAL, True, MONO)]])
txt(s, 8.35, 3.4, 3.9, 0.4,
    [[("Menunjuk ke lokasi CDN video", 9.5, GRAY, False)]])

# catatan penyimpanan
card(s, 8.1, 4.05, 4.3, 2.45, fill=CARD2)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.1, 4.05, 4.3, 0.12, fill=AMBER, radius=0.5)
txt(s, 8.38, 4.32, 3.8, 0.4, [[("PENYIMPANAN VIDEO", 11.5, AMBER, True)]])
storage_notes = [
    "Berkas video TIDAK disimpan di repo Git",
    "(batas GitHub 100 MB/berkas).",
    "",
    "Rencana: Cloudflare R2 sebagai CDN",
    "— 10 GB gratis, sekali unggah.",
]
ny = 4.7
for n in storage_notes:
    if n:
        txt(s, 8.38, ny, 3.8, 0.35, [[(n, 10.5, LIGHT, False)]])
    ny += 0.32
footer(s, 11)
notes(s, """
VIDEO STREAMING. Cara kerja player:
- Komponen VideoPlayer pakai elemen HTML5 <video> bawaan browser.
- Sumber video dirakit dinamis: {BASE}/{id-drama}/{episode}.mp4.
- Fallback otomatis: kalau video gagal dimuat, beralih ke /sample.mp4 + label
  "Video sample — file asli belum diupload". App tetap jalan.
- Navigasi episode prev/next di halaman watch.
Env var NEXT_PUBLIC_VIDEO_BASE_URL menunjuk lokasi video. PENTING: file video
TIDAK ada di repo Git (batas GitHub 100 MB/file) — slide berikutnya menjelaskan
SISTEM penyimpanan video yang sebenarnya dipakai.
""")

# ======================================================================
# SLIDE 12 — Arsitektur Hybrid (Vercel + PC Backup + Tunnel)  [BARU]
# ======================================================================
s = new_slide(BG)
header(s, "Admin, Video & Deployment", "Arsitektur Hybrid", 12)
txt(s, 0.95, 1.88, 11.4, 0.8,
    [[("Aplikasi & video DIPISAH demi hemat biaya — Vercel mengurus tampilan, "
       "PC backup menyalurkan file video besar lewat tunnel.", 13.5, LIGHT, False)]],
    line_spacing=1.3)

# tiga node alur: Pengunjung -> Vercel -> PC Backup
nodes = [
    ("PENGUNJUNG", "Buka app di browser\nHP / desktop", ROSE),
    ("VERCEL", "Aplikasi: halaman, katalog,\nAPI — gratis, selalu online", PURPLE),
    ("PC BACKUP", "Simpan file video asli +\nsalurkan via Caddy & tunnel", TEAL),
]
cw, gx = 3.55, 0.55
x0 = 0.95
ny = 2.95
for i, (t, d, col) in enumerate(nodes):
    x = x0 + i * (cw + gx)
    card(s, x, ny, cw, 1.65, fill=CARD2)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, ny, cw, 0.12, fill=col, radius=0.5)
    txt(s, x + 0.3, ny + 0.32, cw - 0.5, 0.4, [[(t, 15, col, True)]])
    txt(s, x + 0.3, ny + 0.78, cw - 0.55, 0.8,
        [[(ln, 11, LIGHT, False)] for ln in d.split("\n")], line_spacing=1.2)
    if i < 2:
        txt(s, x + cw - 0.04, ny, 0.63, 1.65,
            [[("›", 26, DIM, True)]], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

# alasan + komponen PC backup
card(s, 0.95, 4.95, 5.6, 1.65)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, 4.95, 5.6, 0.12, fill=AMBER, radius=0.5)
txt(s, 1.25, 5.2, 5.0, 0.4, [[("KENAPA DIPISAH?", 12, AMBER, True)]])
txt(s, 1.25, 5.55, 5.1, 1.0,
    [[("File video besar — tak muat di GitHub (batas 100 MB) maupun hosting "
       "gratis. Solusinya: app ringan di Vercel, video tetap di PC. Semua "
       "layanan gratis → biaya Rp 0.", 11, GRAY, False)]], line_spacing=1.25)

card(s, 6.8, 4.95, 5.6, 1.65, fill=CARD2)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.8, 4.95, 5.6, 0.12, fill=PURPLE, radius=0.5)
txt(s, 7.1, 5.2, 5.0, 0.4, [[("KOMPONEN DI PC BACKUP", 12, PURPLE, True)]])
pcbits = [
    ("Caddyfile", "web server penyalur video"),
    ("hardlink-agent.js", "siapkan file video otomatis"),
    ("start-dramaapp.ps1", "1 perintah jalankan semua"),
]
py = 5.55
for name, d in pcbits:
    txt(s, 7.1, py, 5.1, 0.32,
        [[(name, 11, TEAL, True, MONO), ("  —  " + d, 10.5, GRAY, False)]])
    py += 0.34
footer(s, 12)
notes(s, """
ARSITEKTUR HYBRID — INI BAGIAN PALING MENARIK, TONJOLKAN.
Alur: Pengunjung → buka app. App (halaman, katalog, API) jalan di Vercel. Untuk
VIDEO, app mengarah ke PC backup yang menyimpan file asli & menyalurkannya lewat
Caddy + tunnel.
KENAPA begini? File video besar, GitHub batasi 100 MB/file dan hosting gratis tak
muat. Solusi: pisahkan — app ringan di Vercel (gratis), video di PC sendiri.
Hasilnya biaya hosting Rp 0.
Komponen di folder pc-backup-agent/: Caddyfile (server video), hardlink-agent.js
(siapkan video otomatis saat admin tambah drama), start-dramaapp.ps1 (1 perintah
start semua + auto-update Vercel).
Kelemahan jujur yang bisa disebut: URL video sekarang berubah tiap PC restart —
itu yang akan diperbaiki dengan Named Tunnel (lihat slide Perjalanan & Roadmap).
""")

# ======================================================================
# SLIDE 13 — Deployment & DevOps
# ======================================================================
s = new_slide(BG)
header(s, "Admin, Video & Deployment", "Deployment & DevOps", 12)
txt(s, 0.95, 1.85, 11.4, 0.4,
    [[("Dari kode lokal sampai online — alur rilis yang otomatis.",
       12.5, GRAY, False)]])

flow = [
    ("1", "GitHub", "Kode di repositori publik\nmasradenbagus89-ui/dramaapp", ROSE),
    ("2", "Vercel", "Auto-deteksi Next.js,\nbuild & deploy otomatis dari Git", PURPLE),
    ("3", "Online", "Aplikasi live di URL\n.vercel.app — siap diakses", TEAL),
]
cw, gx = 3.55, 0.55
x0 = 0.95
for i, (no, t, d, col) in enumerate(flow):
    x = x0 + i * (cw + gx)
    card(s, x, 2.4, cw, 1.95)
    rect(s, MSO_SHAPE.OVAL, x + 0.28, 2.66, 0.6, 0.6, fill=col)
    txt(s, x + 0.28, 2.66, 0.6, 0.6, [[(no, 18, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 1.05, 2.74, cw - 1.3, 0.5, [[(t, 16, WHITE, True)]])
    txt(s, x + 0.3, 3.4, cw - 0.55, 0.85,
        [[(ln, 10.5, GRAY, False)] for ln in d.split("\n")], line_spacing=1.2)
    if i < 2:
        txt(s, x + cw - 0.05, 2.4, 0.65, 1.95,
            [[("›", 26, DIM, True)]], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

# dua kartu bawah
card(s, 0.95, 4.65, 5.6, 1.95, fill=CARD2)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, 4.65, 5.6, 0.12, fill=AMBER, radius=0.5)
txt(s, 1.25, 4.9, 5.0, 0.4, [[("OTOMASI", 12, AMBER, True)]])
txt(s, 1.25, 5.22, 5.05, 1.3,
    [[("Skrip ", 11.5, LIGHT, False),
      ("start-dramaapp.ps1", 11.5, AMBER, True, MONO)],
     [("Satu perintah menjalankan seluruh", 11.5, GRAY, False)],
     [("lingkungan dev + sinkron otomatis ke Vercel.", 11.5, GRAY, False)]],
    line_spacing=1.3, space_after=3)

card(s, 6.8, 4.65, 5.6, 1.95)
rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.8, 4.65, 5.6, 0.12, fill=PURPLE, radius=0.5)
txt(s, 7.1, 4.9, 5.0, 0.4, [[("KONFIGURASI", 12, PURPLE, True)]])
txt(s, 7.1, 5.22, 5.1, 1.3,
    [[("5 environment variable", 11.5, LIGHT, True)],
     [("diatur di dashboard Vercel — termasuk", 11.5, GRAY, False)],
     [("URL basis video & pengaturan aplikasi.", 11.5, GRAY, False)]],
    line_spacing=1.3, space_after=3)
footer(s, 12)
notes(s, """
DEPLOYMENT & DEVOPS. Alur rilis otomatis: (1) Kode di GitHub (repo publik
masradenbagus89-ui/dramaapp) -> (2) Vercel auto-deteksi Next.js, build & deploy
otomatis tiap ada perubahan -> (3) Online di URL .vercel.app, siap diakses.
OTOMASI: skrip start-dramaapp.ps1 menjalankan seluruh lingkungan dev + sinkron
otomatis ke Vercel dengan satu perintah.
KONFIGURASI: 5 environment variable di dashboard Vercel (termasuk URL basis video).
Pesan: cukup push ke GitHub -> otomatis online, tanpa langkah manual.
""")

# ======================================================================
# SLIDE 14 — Perjalanan Pembuatan (8 Tahap)  [BARU]
# ======================================================================
s = new_slide(BG)
header(s, "Pembelajaran", "Perjalanan Pembuatan Project", 14)
txt(s, 0.95, 1.85, 11.4, 0.4,
    [[("8 tahap dari nol sampai online — 6 sudah selesai, 2 berikutnya.",
       12.5, GRAY, False)]])

stages = [
    ("1", "Setup Project", "Ambil code dari GitHub, jalan di laptop", True),
    ("2", "Online di Vercel", "App bisa dibuka publik: dramaapp.vercel.app", True),
    ("3", "Video Bisa Diputar", "Video dari PC backup disalurkan via tunnel", True),
    ("4", "Admin Form Otomatis", "Tambah drama tanpa edit file manual", True),
    ("5", "Auto-Hardlink", "Tambah drama baru cukup ~3 klik / 2 menit", True),
    ("6", "Script Otomatis", "1 command start semua + auto-update Vercel", True),
    ("7", "Named Tunnel", "Berikutnya: URL video stabil selamanya", False),
    ("8", "Perbaikan Akhir", "Rencana: fix bug judul panjang + keamanan", False),
]
cw, ch, gx, gy = 3.69, 0.96, 0.19, 0.2
x0, y0 = 0.95, 2.4
for i, (no, t, d, done) in enumerate(stages):
    r, c = divmod(i, 2)
    x = x0 + c * (cw + gx)
    # kolom 3 (tahap 7-8) jadi baris ke-4; pakai 2 kolom x 4 baris
    x = 0.95 + c * (5.78 + 0.3)
    y = y0 + r * (ch + gy)
    col = TEAL if done else AMBER
    bg = CARD if done else CARD2
    card(s, x, y, 5.78, ch, fill=bg)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.12, ch, fill=col, radius=0.5)
    rect(s, MSO_SHAPE.OVAL, x + 0.3, y + 0.27, 0.42, 0.42, fill=col)
    txt(s, x + 0.3, y + 0.27, 0.42, 0.42, [[(no, 15, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    mark = "SELESAI" if done else ("BERIKUTNYA" if no == "7" else "RENCANA")
    txt(s, x + 0.9, y + 0.16, 4.7, 0.35,
        [[(t, 13.5, WHITE, True), ("   " + mark, 8.5, col, True)]])
    txt(s, x + 0.9, y + 0.5, 4.75, 0.4, [[(d, 10, GRAY, False)]], line_spacing=1.1)
footer(s, 14)
notes(s, """
PERJALANAN PEMBUATAN. Ceritakan ini sebagai "SISTEM yang dipakai selama membangun".
8 tahap, 6 SELESAI: (1) Setup project dari GitHub, (2) Online di Vercel,
(3) Video bisa diputar lewat PC backup + tunnel, (4) Admin form otomatis (tak
perlu edit file manual), (5) Auto-hardlink (tambah drama cukup ~3 klik / 2 menit),
(6) Script otomatis (1 command start semua + auto-update Vercel).
2 BERIKUTNYA: (7) Named Tunnel — biar URL video stabil selamanya, berhenti
ganti-ganti tiap PC restart; (8) Perbaikan akhir — fix bug judul drama panjang
(slug terpotong) + perkuat keamanan login admin.
Pesan: tunjukkan progres nyata & produktivitas yang makin meningkat tiap tahap.
""")

# ======================================================================
# SLIDE 15 — Tantangan & Solusi
# ======================================================================
s = new_slide(BG)
header(s, "Pembelajaran", "Tantangan Teknis & Solusi", 13)
ch = [
    ("Jaringan server terkunci",
     "Firewall memblokir port masuk; lalu lintas keluar dipaksa lewat proxy yang "
     "memblokir layanan tunnel (ngrok, cloudflared, localtunnel — semua gagal).",
     "Deploy ke cloud (Vercel) — aplikasi diakses lewat URL publik tanpa "
     "bergantung pada jaringan server."),
    ("Berkas video terlalu besar",
     "Video melebihi batas 100 MB GitHub sehingga tidak bisa ikut masuk repo.",
     "Pisahkan video ke penyimpanan eksternal (Cloudflare R2 / CDN); path "
     "diatur lewat environment variable."),
    ("Tanpa database",
     "Belum ada server basis data untuk menyimpan drama, komentar, dan interaksi.",
     "Gunakan penyimpanan berkas JSON — sederhana, cukup untuk skala awal, "
     "mudah ditingkatkan ke database nanti."),
]
y = 2.0
for title_t, prob, sol in ch:
    card(s, 0.95, y, 11.45, 1.45)
    txt(s, 1.25, y + 0.16, 4.0, 1.2, [[(title_t, 14, WHITE, True)]],
        line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 4.55, y + 0.18, 0.09, 1.1,
         fill=LINEC, radius=0.5)
    txt(s, 4.8, y + 0.17, 3.55, 0.3, [[("TANTANGAN", 9, ROSE, True)]])
    txt(s, 4.8, y + 0.43, 3.6, 1.0,
        [[(prob, 9.7, GRAY, False)]], line_spacing=1.2)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.55, y + 0.18, 0.09, 1.1,
         fill=TEAL, radius=0.5)
    txt(s, 8.8, y + 0.17, 3.5, 0.3, [[("SOLUSI", 9, TEAL, True)]])
    txt(s, 8.8, y + 0.43, 3.45, 1.0,
        [[(sol, 9.7, LIGHT, False)]], line_spacing=1.2)
    y += 1.57
footer(s, 13)
notes(s, """
TANTANGAN & SOLUSI. Cerita problem nyata -> solusi (bikin presentasi hidup):
1. Jaringan server terkunci — firewall blokir port masuk, proxy blokir tunnel
   (ngrok, cloudflared, localtunnel semua gagal). Solusi: deploy ke Vercel,
   diakses lewat URL publik tanpa bergantung jaringan server.
2. Video terlalu besar — lewat batas 100 MB GitHub. Solusi: pisahkan video ke
   penyimpanan eksternal/PC backup; path diatur lewat env var.
3. Tanpa database — belum ada server DB. Solusi: pakai file JSON, sederhana,
   cukup untuk skala awal, mudah di-upgrade ke database nanti.
Pesan: menghadapi kendala teknis nyata dan menyelesaikannya dengan keputusan
arsitektur yang tepat.
""")

# ======================================================================
# SLIDE 14 — Roadmap
# ======================================================================
s = new_slide(BG)
header(s, "Langkah Berikutnya", "Roadmap Pengembangan", 14)

txt(s, 0.95, 1.95, 5.6, 0.4, [[("JANGKA PENDEK", 12, ROSE, True)]])
short = [
    "Unggah video ke Cloudflare R2 & perbarui path",
    "Uji seluruh halaman setelah deploy",
    "Identifikasi & perbaiki bug satu per satu",
]
y = 2.4
for it in short:
    card(s, 0.95, y, 5.6, 0.78)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, y, 0.1, 0.78, fill=ROSE, radius=0.5)
    txt(s, 1.3, y, 5.1, 0.78, [[(it, 11.5, LIGHT, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y += 0.9

txt(s, 6.85, 1.95, 5.5, 0.4, [[("JANGKA PANJANG", 12, PURPLE, True)]])
long = [
    "Tingkatkan dari JSON store ke database",
    "Autentikasi lebih kuat (kini berbasis localStorage)",
    "Fitur pencarian & rekomendasi konten",
    "Optimasi performa & infinite scroll Shorts",
]
y = 2.4
for it in long:
    card(s, 6.85, y, 5.55, 0.78, fill=CARD2)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.85, y, 0.1, 0.78, fill=PURPLE, radius=0.5)
    txt(s, 7.2, y, 5.05, 0.78, [[(it, 11.5, LIGHT, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y += 0.9

card(s, 0.95, 5.55, 11.45, 1.0, fill=RGBColor(0x10, 0x12, 0x1C))
txt(s, 1.3, 5.55, 11.0, 1.0,
    [[("Status saat ini:  ", 12, GRAY, True),
      ("aplikasi fungsional & ter-deploy — fokus berikutnya adalah "
       "memutar video dari CDN dan stabilisasi.", 12, LIGHT, False)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
footer(s, 14)
notes(s, """
ROADMAP PENGEMBANGAN. JANGKA PENDEK: unggah video ke Cloudflare R2 & perbarui
path, uji seluruh halaman setelah deploy, perbaiki bug satu per satu.
JANGKA PANJANG: upgrade dari JSON store ke database, autentikasi lebih kuat
(kini berbasis localStorage), fitur pencarian & rekomendasi, optimasi performa &
infinite scroll Shorts.
Status sekarang: aplikasi fungsional & ter-deploy — fokus berikutnya memutar
video dari CDN dan stabilisasi (Named Tunnel).
""")

# ======================================================================
# SLIDE 17 — Antisipasi Tanya-Jawab  [BARU]
# ======================================================================
s = new_slide(BG)
header(s, "Pembelajaran", "Antisipasi Tanya-Jawab", 17)
txt(s, 0.95, 1.85, 11.4, 0.4,
    [[("Pertanyaan yang mungkin muncul — siapkan jawaban singkat ini.",
       12.5, GRAY, False)]])

qa = [
    ("Kenapa tidak pakai database beneran?",
     "File JSON cukup untuk skala ini, lebih ringan & tanpa biaya. Kalau data "
     "besar, tinggal migrasi ke database tanpa ubah tampilan.", ROSE),
    ("Berapa biaya hosting?",
     "Rp 0 — Vercel free tier untuk aplikasi, PC sendiri untuk video, semua "
     "tool open-source.", AMBER),
    ("Apakah aman?",
     "Role admin diverifikasi di server di tiap endpoint (bukan cuma tampilan); "
     "video dipisah dari aplikasi.", TEAL),
    ("Bisa diakses dari HP?",
     "Bisa. Desain responsive, ada bottom-nav khusus mobile, jalan di browser "
     "HP mana saja.", PURPLE),
]
cw, ch, gx, gy = 5.6, 1.95, 0.25, 0.22
x0, y0 = 0.95, 2.4
for i, (q, a, col) in enumerate(qa):
    r, c = divmod(i, 2)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    card(s, x, y, cw, ch)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, 0.12, fill=col, radius=0.5)
    txt(s, x + 0.3, y + 0.32, cw - 0.55, 0.5,
        [[("T:  ", 13, col, True), (q, 13, WHITE, True)]], line_spacing=1.1)
    txt(s, x + 0.3, y + 0.95, cw - 0.58, 0.95,
        [[("J:  ", 11.5, col, True), (a, 11.5, GRAY, False)]], line_spacing=1.25)
footer(s, 17)
notes(s, """
ANTISIPASI TANYA-JAWAB. Hafalkan jawaban singkat 4 pertanyaan ini:
1. Kenapa tanpa database? -> JSON cukup untuk skala ini, ringan, gratis; mudah
   migrasi ke DB nanti tanpa ubah tampilan.
2. Berapa biaya? -> Rp 0: Vercel free tier + PC sendiri + tool open-source.
3. Aman? -> Role admin diverifikasi di server tiap endpoint; video terpisah.
4. Bisa di HP? -> Bisa, responsive + bottom-nav mobile, jalan di browser HP apa pun.
Jawab tenang dan singkat. Kalau tak tahu, jujur "itu masuk rencana pengembangan".
""")

# ======================================================================
# SLIDE 18 — Penutup
# ======================================================================
s = new_slide(BG)
rect(s, MSO_SHAPE.OVAL, -2.4, 3.6, 6.4, 6.4, fill=GLOW_P)
rect(s, MSO_SHAPE.OVAL, 9.6, -2.6, 6.0, 6.0, fill=GLOW_R)
rect(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16, fill=ROSE)

txt(s, 1.0, 1.05, 11.0, 0.4, [[("RINGKASAN", 13, PURPLE, True)]])
txt(s, 0.97, 1.4, 11.4, 1.0,
    [[("DramaKu", 46, WHITE, True),
      (" — modern, ringan, siap berkembang", 30, LIGHT, False)]])

summary = [
    ("Modern", "Next.js 16, React 19 & TypeScript dalam satu codebase full-stack.", ROSE),
    ("Mobile-first", "Pengalaman menonton dirancang utama untuk layar HP.", PURPLE),
    ("Praktis", "Data berbasis file & panel admin — cepat dikelola tanpa kode.", AMBER),
    ("Ter-deploy", "Online di Vercel dengan alur rilis otomatis dari GitHub.", TEAL),
]
cw, gx = 2.69, 0.27
x0 = 0.95
for i, (t, d, col) in enumerate(summary):
    x = x0 + i * (cw + gx)
    card(s, x, 2.7, cw, 1.95)
    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.26, 2.96, 0.5, 0.11,
         fill=col, radius=0.5)
    txt(s, x + 0.26, 3.18, cw - 0.5, 0.4, [[(t, 14, WHITE, True)]])
    txt(s, x + 0.26, 3.58, cw - 0.5, 1.0,
        [[(d, 10.5, GRAY, False)]], line_spacing=1.25)

grad(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, 5.1, 11.45, 1.35, CARD2, CARD,
     angle=0, radius=0.08)
txt(s, 1.35, 5.3, 7.5, 0.5, [[("Terima kasih", 24, WHITE, True)]])
txt(s, 1.35, 5.85, 8.0, 0.4,
    [[("Siap untuk sesi tanya-jawab & demo aplikasi", 12.5, GRAY, False)]])
txt(s, 8.7, 5.42, 3.5, 0.4,
    [[("REPOSITORI", 10, PURPLE, True)]], align=PP_ALIGN.RIGHT)
txt(s, 7.7, 5.7, 4.5, 0.4,
    [[("github.com/masradenbagus89-ui/dramaapp", 11, AMBER, True, MONO)]],
    align=PP_ALIGN.RIGHT)

txt(s, 1.0, 6.85, 11.3, 0.35,
    [[("DramaKu  ·  Gambaran Teknis Project  ·  2 Juni 2026", 9.5, DIM, False)]])
notes(s, """
PENUTUP. Kalimat penutup:
"Singkatnya, DramaApp adalah aplikasi streaming yang sudah online, dibangun
dengan teknologi modern (Next.js, TypeScript, Vercel), memakai arsitektur hybrid
yang hemat biaya untuk menangani video besar, dan punya sistem admin yang membuat
penambahan konten jadi sangat cepat. Dari 8 tahap rencana, 6 sudah selesai dan
aplikasinya sudah bisa dipakai sekarang."
Lalu buka sesi tanya-jawab & tawarkan demo aplikasi langsung.
Angka penutup: 33 judul, 7 kategori, 6/8 tahap selesai, biaya Rp 0,
URL dramaapp.vercel.app.
""")

# ----------------------------------------------------------------------
out = r"D:\Users\user18\dramaapp\presentasi\DramaApp-Presentasi-Teknis.pptx"
prs.save(out)
print("OK ->", out)
print("Total slide:", len(prs.slides._sldIdLst))
